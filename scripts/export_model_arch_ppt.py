from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = PROJECT_ROOT / "outputs"
OUTPUT_PATH = OUTPUT_DIR / "UNetResNet34_architecture.pptx"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Calibri"

BG = RGBColor(247, 243, 235)
PANEL = RGBColor(255, 252, 246)
TITLE_BAR = RGBColor(32, 58, 84)
TITLE_TEXT = RGBColor(248, 247, 244)
TEXT_DARK = RGBColor(34, 34, 34)
TEXT_MUTED = RGBColor(85, 85, 85)
LINE = RGBColor(56, 60, 66)
SHADOW = RGBColor(212, 203, 191)

INPUT = RGBColor(54, 102, 135)
ENCODER = RGBColor(76, 136, 153)
POOL = RGBColor(119, 126, 140)
CENTER = RGBColor(53, 120, 78)
DECODER = RGBColor(204, 138, 60)
HEAD = RGBColor(128, 83, 146)
OUTPUT = RGBColor(172, 70, 86)
SKIP = RGBColor(190, 84, 76)
ACCENT = RGBColor(216, 168, 92)


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_full_rect(slide, color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        SLIDE_W,
        SLIDE_H,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.78),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_BAR
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(8.8), Inches(0.45))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.name = FONT_CN
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = TITLE_TEXT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = FONT_CN
        r2.font.size = Pt(10)
        r2.font.color.rgb = rgb(219, 225, 230)


def add_textbox(slide, left, top, width, height, text, font_size=12, color=TEXT_DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT_CN
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_shadow_card(slide, left, top, width, height, title, body, fill, title_size=13, body_size=10, text_color=TITLE_TEXT):
    shadow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left + Inches(0.06),
        top + Inches(0.06),
        width,
        height,
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = SHADOW
    shadow.line.fill.background()

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = LINE
    card.line.width = Pt(1.2)

    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = FONT_CN
    r1.font.size = Pt(title_size)
    r1.font.bold = True
    r1.font.color.rgb = text_color

    if body:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = FONT_EN
        r2.font.size = Pt(body_size)
        r2.font.color.rgb = text_color

    return card


def add_panel(slide, left, top, width, height, title=None):
    shadow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left + Inches(0.06),
        top + Inches(0.06),
        width,
        height,
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = SHADOW
    shadow.line.fill.background()

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = rgb(199, 190, 177)
    panel.line.width = Pt(1.0)

    if title:
        add_textbox(
            slide,
            left + Inches(0.18),
            top + Inches(0.08),
            width - Inches(0.36),
            Inches(0.25),
            title,
            font_size=11,
            color=TEXT_MUTED,
            bold=True,
        )
    return panel


def add_chip(slide, left, top, width, height, text, fill, color=TEXT_DARK):
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    chip.line.fill.background()
    tf = chip.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_CN
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = color
    return chip


def add_arrow(slide, x1, y1, x2, y2, color=LINE, width=1.8, dashed=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass
    return line


def mid_left(shape):
    return shape.left, shape.top + shape.height / 2


def mid_right(shape):
    return shape.left + shape.width, shape.top + shape.height / 2


def mid_top(shape):
    return shape.left + shape.width / 2, shape.top


def mid_bottom(shape):
    return shape.left + shape.width / 2, shape.top + shape.height


def build_cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, BG)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(241, 233, 219)
    accent.line.fill.background()

    wave = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON,
        Inches(8.85),
        Inches(0.8),
        Inches(4.0),
        Inches(5.6),
    )
    wave.fill.solid()
    wave.fill.fore_color.rgb = rgb(230, 214, 192)
    wave.line.fill.background()

    title_card = add_panel(slide, Inches(0.6), Inches(1.0), Inches(7.7), Inches(2.55))
    add_textbox(slide, Inches(0.9), Inches(1.35), Inches(6.9), Inches(0.55), "UNetResNet34 模型结构图", 28, TITLE_BAR, True)
    add_textbox(
        slide,
        Inches(0.92),
        Inches(2.00),
        Inches(6.8),
        Inches(0.95),
        "基于你当前 `src/model.py` 的真实实现导出\n"
        "包含编码器、瓶颈层、解码器、skip connection 与 DecoderBlock 细化流程",
        13,
        TEXT_MUTED,
    )

    add_chip(slide, Inches(0.95), Inches(3.00), Inches(1.05), Inches(0.28), "Encoder", rgb(212, 231, 236))
    add_chip(slide, Inches(2.10), Inches(3.00), Inches(1.05), Inches(0.28), "Bottleneck", rgb(214, 233, 219))
    add_chip(slide, Inches(3.30), Inches(3.00), Inches(1.05), Inches(0.28), "Decoder", rgb(244, 224, 194))
    add_chip(slide, Inches(4.45), Inches(3.00), Inches(1.15), Inches(0.28), "Seg Head", rgb(227, 213, 235))

    add_panel(slide, Inches(8.15), Inches(1.15), Inches(4.45), Inches(4.55), "内容概览")
    items = [
        ("01", "整体架构页", "把 ResNet34 编码器、U-Net 解码器与 skip 连接放在一页里。"),
        ("02", "尺寸与通道流", "按 forward 的顺序列出 x0 / x1 / x2 / x3 / x4 / d4 / d3 / d2 / d1。"),
        ("03", "DecoderBlock 细节", "单独解释 interpolate、cat(dim=1)、ConvBlock 是怎么连起来的。"),
    ]
    base_top = 1.55
    for idx, (num, title, desc) in enumerate(items):
        y = Inches(base_top + idx * 1.23)
        add_chip(slide, Inches(8.45), y, Inches(0.52), Inches(0.3), num, ACCENT, TITLE_BAR)
        add_textbox(slide, Inches(9.08), y - Inches(0.02), Inches(2.9), Inches(0.25), title, 13, TITLE_BAR, True)
        add_textbox(slide, Inches(9.08), y + Inches(0.28), Inches(3.0), Inches(0.52), desc, 10, TEXT_MUTED)

    add_textbox(slide, Inches(0.75), Inches(6.72), Inches(4.5), Inches(0.25), "生成文件: outputs/UNetResNet34_architecture.pptx", 10, TEXT_MUTED)


def build_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, BG)
    add_title_bar(slide, "整体架构", "按你当前 model.py 的 forward 路径整理，不是抽象示意图。")

    add_panel(slide, Inches(0.25), Inches(0.95), Inches(3.75), Inches(5.95), "Encoder / ResNet34")
    add_panel(slide, Inches(4.18), Inches(0.95), Inches(2.65), Inches(5.95), "Bridge")
    add_panel(slide, Inches(7.00), Inches(0.95), Inches(4.55), Inches(5.95), "Decoder / U-Net")
    add_panel(slide, Inches(11.73), Inches(0.95), Inches(1.35), Inches(5.95), "Output")

    left = Inches(0.63)
    width = Inches(3.0)
    height = Inches(0.63)

    input_box = add_shadow_card(slide, left, Inches(1.16), width, height, "输入图像", "3 x H x W", INPUT)
    stem = add_shadow_card(slide, left, Inches(2.00), width, height, "encoder_stem", "conv1 / bn1 / relu\n64 @ H/2", ENCODER)
    pool = add_shadow_card(slide, left, Inches(2.78), width, Inches(0.46), "maxpool", "64 @ H/4", POOL, 12, 10)
    x1 = add_shadow_card(slide, left, Inches(3.40), width, height, "encoder_layer1  -> x1", "64 @ H/4", ENCODER)
    x2 = add_shadow_card(slide, left, Inches(4.20), width, height, "encoder_layer2  -> x2", "128 @ H/8", ENCODER)
    x3 = add_shadow_card(slide, left, Inches(5.00), width, height, "encoder_layer3  -> x3", "256 @ H/16", ENCODER)
    x4 = add_shadow_card(slide, left, Inches(5.80), width, height, "encoder_layer4  -> x4", "512 @ H/32", ENCODER)

    center = add_shadow_card(slide, Inches(4.52), Inches(3.55), Inches(1.95), Inches(0.95), "center", "ConvBlock\n512 -> 512", CENTER)

    d4 = add_shadow_card(slide, Inches(7.28), Inches(5.00), Inches(3.98), height, "decoder4", "x4 + skip(x3)\n512 + 256 -> 256 @ H/16", DECODER, 13, 10)
    d3 = add_shadow_card(slide, Inches(7.28), Inches(4.20), Inches(3.98), height, "decoder3", "d4 + skip(x2)\n256 + 128 -> 128 @ H/8", DECODER, 13, 10)
    d2 = add_shadow_card(slide, Inches(7.28), Inches(3.40), Inches(3.98), height, "decoder2", "d3 + skip(x1)\n128 + 64 -> 64 @ H/4", DECODER, 13, 10)
    d1 = add_shadow_card(slide, Inches(7.28), Inches(2.60), Inches(3.98), height, "decoder1", "d2 + skip(x0)\n64 + 64 -> 64 @ H/2", DECODER, 13, 10)
    fin = add_shadow_card(slide, Inches(7.28), Inches(1.75), Inches(3.98), Inches(0.56), "final interpolate", "64 @ H x W", DECODER, 12, 10)

    head = add_shadow_card(slide, Inches(11.92), Inches(2.95), Inches(0.95), Inches(1.05), "segmentation\nhead", "64 -> 32 -> 1", HEAD, 12, 10)
    out = add_shadow_card(slide, Inches(11.92), Inches(4.30), Inches(0.95), Inches(0.8), "输出", "1 x H x W", OUTPUT, 12, 10)

    for src, dst in [
        (input_box, stem),
        (stem, pool),
        (pool, x1),
        (x1, x2),
        (x2, x3),
        (x3, x4),
    ]:
        add_arrow(slide, *mid_bottom(src), *mid_top(dst))

    add_arrow(slide, *mid_right(x4), *mid_left(center))
    add_arrow(slide, *mid_right(center), *mid_left(d4))
    add_arrow(slide, *mid_top(d4), *mid_bottom(d3))
    add_arrow(slide, *mid_top(d3), *mid_bottom(d2))
    add_arrow(slide, *mid_top(d2), *mid_bottom(d1))
    add_arrow(slide, *mid_top(d1), *mid_bottom(fin))
    add_arrow(slide, *mid_right(fin), *mid_left(head))
    add_arrow(slide, *mid_bottom(head), *mid_top(out))

    skip_pairs = [
        (stem, d1, "skip x0"),
        (x1, d2, "skip x1"),
        (x2, d3, "skip x2"),
        (x3, d4, "skip x3"),
    ]
    for src, dst, label in skip_pairs:
        x1p, y1p = mid_right(src)
        x2p, y2p = mid_left(dst)
        add_arrow(slide, x1p, y1p, x2p, y2p, color=SKIP, width=1.5, dashed=True)
        add_textbox(
            slide,
            (x1p + x2p) / 2 - Inches(0.42),
            y2p - Inches(0.12),
            Inches(0.85),
            Inches(0.2),
            label,
            9,
            SKIP,
            True,
            PP_ALIGN.CENTER,
        )

    add_chip(slide, Inches(0.6), Inches(6.28), Inches(1.6), Inches(0.28), "实线 = 主分支前向", rgb(227, 232, 236))
    add_chip(slide, Inches(2.35), Inches(6.28), Inches(1.9), Inches(0.28), "虚线 = skip connection", rgb(248, 224, 222))
    add_chip(slide, Inches(4.45), Inches(6.28), Inches(1.8), Inches(0.28), "颜色 = 模块类型", rgb(240, 232, 219))

    add_textbox(
        slide,
        Inches(6.55),
        Inches(6.18),
        Inches(6.1),
        Inches(0.48),
        "forward: x0 = stem(x) -> x1 = layer1(pool(x0)) -> x2 -> x3 -> x4 -> center -> d4 -> d3 -> d2 -> d1 -> head",
        10,
        TEXT_MUTED,
    )


def build_shape_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, BG)
    add_title_bar(slide, "尺寸与通道流", "以下按输入 640 x 640 举例；如果输入换成 384，比例关系不变。")

    add_panel(slide, Inches(0.35), Inches(1.0), Inches(12.6), Inches(5.95), "forward() 中各张量的尺寸变化")

    boxes = []
    specs = [
        ("x", "3 x 640 x 640", INPUT),
        ("x0", "64 x 320 x 320", ENCODER),
        ("x1", "64 x 160 x 160", ENCODER),
        ("x2", "128 x 80 x 80", ENCODER),
        ("x3", "256 x 40 x 40", ENCODER),
        ("x4", "512 x 20 x 20", ENCODER),
        ("center(x4)", "512 x 20 x 20", CENTER),
        ("d4", "256 x 40 x 40", DECODER),
        ("d3", "128 x 80 x 80", DECODER),
        ("d2", "64 x 160 x 160", DECODER),
        ("d1", "64 x 320 x 320", DECODER),
        ("logits", "1 x 640 x 640", OUTPUT),
    ]

    top_row_y = Inches(2.0)
    bot_row_y = Inches(4.15)
    x_positions_top = [Inches(0.65), Inches(2.05), Inches(3.45), Inches(4.85), Inches(6.25), Inches(7.65)]
    x_positions_bot = [Inches(7.65), Inches(6.15), Inches(4.65), Inches(3.15), Inches(1.65), Inches(0.45)]
    size_w = Inches(1.1)
    size_h = Inches(0.82)

    for idx, (name, spec, color) in enumerate(specs[:6]):
        card = add_shadow_card(slide, x_positions_top[idx], top_row_y, size_w, size_h, name, spec, color, 12, 10)
        boxes.append(card)
    for idx, (name, spec, color) in enumerate(specs[6:]):
        card = add_shadow_card(slide, x_positions_bot[idx], bot_row_y, size_w, size_h, name, spec, color, 12, 10)
        boxes.append(card)

    for i in range(5):
        add_arrow(slide, *mid_right(boxes[i]), *mid_left(boxes[i + 1]))
    add_arrow(slide, *mid_bottom(boxes[5]), *mid_top(boxes[6]))
    for i in range(6, 11):
        add_arrow(slide, *mid_left(boxes[i]), *mid_right(boxes[i + 1]))

    add_textbox(slide, Inches(2.3), Inches(1.56), Inches(1.0), Inches(0.2), "stem", 10, TEXT_MUTED, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.68), Inches(1.56), Inches(1.0), Inches(0.2), "pool + layer1", 10, TEXT_MUTED, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.11), Inches(1.56), Inches(0.8), Inches(0.2), "layer2", 10, TEXT_MUTED, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.53), Inches(1.56), Inches(0.8), Inches(0.2), "layer3", 10, TEXT_MUTED, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.94), Inches(1.56), Inches(0.8), Inches(0.2), "layer4", 10, TEXT_MUTED, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.13), Inches(3.52), Inches(1.0), Inches(0.25), "ConvBlock", 10, TEXT_MUTED, True, PP_ALIGN.CENTER)

    add_note = add_panel(slide, Inches(9.0), Inches(1.55), Inches(3.55), Inches(4.6), "关键观察")
    del add_note
    notes = [
        "1. `encoder_stem` 只包含 conv1 / bn1 / relu，所以 x0 是 H/2。",
        "2. `encoder_pool` 发生在进入 layer1 之前，所以 x1 是 H/4。",
        "3. `center` 不改分辨率，只在最深层做 512 -> 512 的卷积融合。",
        "4. `decoder4/3/2/1` 每次都先上采样，再和 skip 特征按通道拼接。",
        "5. `d1` 仍然是 H/2，所以 forward 里额外做一次 final interpolate 回到输入大小。",
    ]
    for idx, text in enumerate(notes):
        add_textbox(slide, Inches(9.22), Inches(1.9 + idx * 0.76), Inches(3.0), Inches(0.52), text, 10, TEXT_DARK)

    add_chip(slide, Inches(0.75), Inches(6.42), Inches(1.35), Inches(0.28), "通道 x 高 x 宽", rgb(234, 229, 222))


def build_decoder_detail_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_rect(slide, BG)
    add_title_bar(slide, "DecoderBlock 细化", "这一页专门解释：为什么先拼接，再卷积；为什么 concat 后通道会变化。")

    add_panel(slide, Inches(0.4), Inches(1.0), Inches(12.45), Inches(5.95), "DecoderBlock.forward(x, skip)")

    x_box = add_shadow_card(slide, Inches(0.8), Inches(2.05), Inches(2.1), Inches(0.78), "x", "B x in_channels x h x w", DECODER, 14, 11)
    skip_box = add_shadow_card(slide, Inches(0.8), Inches(4.30), Inches(2.1), Inches(0.78), "skip", "B x skip_channels x H x W", ENCODER, 14, 11)
    interp = add_shadow_card(slide, Inches(3.55), Inches(2.05), Inches(2.2), Inches(0.96), "F.interpolate", "size = skip.shape[-2:]\n双线性上采样", DECODER, 14, 10)
    cat = add_shadow_card(slide, Inches(6.35), Inches(3.05), Inches(2.4), Inches(0.96), "torch.cat([x, skip], dim=1)", "B x (in + skip) x H x W", CENTER, 12, 10)
    conv = add_shadow_card(slide, Inches(9.15), Inches(3.05), Inches(2.05), Inches(0.96), "ConvBlock", "ConvBlock(in + skip, out)", HEAD, 14, 10)
    out = add_shadow_card(slide, Inches(11.55), Inches(3.05), Inches(0.9), Inches(0.96), "out", "B x out x H x W", OUTPUT, 14, 10)

    add_arrow(slide, *mid_right(x_box), *mid_left(interp))
    add_arrow(slide, *mid_right(interp), *mid_left(cat))
    add_arrow(slide, *mid_right(skip_box), *mid_left(cat), color=SKIP, width=1.5, dashed=True)
    add_arrow(slide, *mid_right(cat), *mid_left(conv))
    add_arrow(slide, *mid_right(conv), *mid_left(out))

    add_textbox(slide, Inches(3.85), Inches(1.55), Inches(1.5), Inches(0.24), "先对齐空间尺寸", 10, TEXT_MUTED, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.75), Inches(2.55), Inches(1.55), Inches(0.24), "再按通道拼接", 10, TEXT_MUTED, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.35), Inches(2.55), Inches(1.55), Inches(0.24), "最后卷积融合", 10, TEXT_MUTED, True, PP_ALIGN.CENTER)

    add_panel(slide, Inches(0.9), Inches(5.55), Inches(3.6), Inches(0.9), "例子")
    add_textbox(
        slide,
        Inches(1.15),
        Inches(5.82),
        Inches(3.1),
        Inches(0.38),
        "Decoder3: in=256, skip=128, out=128\ncat 之后是 384 通道，再由 ConvBlock 压回 128。",
        11,
        TEXT_DARK,
    )

    add_panel(slide, Inches(4.85), Inches(5.55), Inches(3.2), Inches(0.9), "对应代码")
    add_textbox(
        slide,
        Inches(5.08),
        Inches(5.82),
        Inches(2.72),
        Inches(0.34),
        "x = torch.cat([x, skip], dim=1)\n"
        "x = self.conv_block(x)",
        11,
        TEXT_DARK,
    )

    add_panel(slide, Inches(8.35), Inches(5.55), Inches(4.0), Inches(0.9), "容易混淆的点")
    add_textbox(
        slide,
        Inches(8.58),
        Inches(5.82),
        Inches(3.5),
        Inches(0.36),
        "`in_channels` 是主分支 x 的通道数。\n`skip_channels` 是 skip 特征的通道数，不是拼接后的总通道数。",
        11,
        TEXT_DARK,
    )


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover_slide(prs)
    build_overview_slide(prs)
    build_shape_flow_slide(prs)
    build_decoder_detail_slide(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation()
    print(OUTPUT_PATH)
